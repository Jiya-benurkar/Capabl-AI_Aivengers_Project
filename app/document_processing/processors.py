"""
Document Processors
-------------------
Handles PDF, DOCX, and PPTX files.
Extracts text with metadata and auto-classifies content type.
"""

from __future__ import annotations

import re
import os
from pathlib import Path
from typing import List, Dict, Any

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.config import CHUNK_SIZE, CHUNK_OVERLAP


# ─────────────────────────────────────────────────────────────────
#  Content Type Classifier
# ─────────────────────────────────────────────────────────────────

_QUESTION_PAPER_SIGNALS = [
    r"\bQ\.\s*\d+", r"\bQuestion\s+\d+", r"\bmarks?\b",
    r"\b(20\d{2}|19\d{2})\s*(paper|exam|question)", r"\banswer\s+any\b",
    r"\bmax(imum)?\s+marks\b", r"\btime\s*:\s*\d+\s*h(ours?)?\b",
]

_LAB_MANUAL_SIGNALS = [
    r"\bexperiment\s+no\.?\s*\d+", r"\baim\s*:", r"\bprocedure\s*:",
    r"\bobservation\s*:", r"\bresult\s*:", r"\bviva\s+questions?\b",
    r"\bpre\s*-?\s*lab\b",
]

_LECTURE_NOTES_SIGNALS = [
    r"\blecture\s+\d+", r"\bunit\s+[ivxIVX\d]+", r"\btopic\s*:", 
    r"\blearning\s+objectives?\b", r"^[-•●]\s+",
]


def classify_content_type(text: str, filename: str) -> str:
    """Heuristically classify a document as one of:
       question_paper | lab_manual | lecture_notes | textbook
    """
    sample = text[:3000].lower()
    fname = filename.lower()

    # Filename clues
    if any(k in fname for k in ["qp", "question", "paper", "exam", "past"]):
        return "question_paper"
    if any(k in fname for k in ["lab", "practical", "experiment"]):
        return "lab_manual"
    if any(k in fname for k in ["note", "lecture", "slide", "ppt"]):
        return "lecture_notes"

    # Content signals
    qp_hits = sum(1 for p in _QUESTION_PAPER_SIGNALS if re.search(p, sample, re.I))
    lab_hits = sum(1 for p in _LAB_MANUAL_SIGNALS if re.search(p, sample, re.I))
    lec_hits = sum(1 for p in _LECTURE_NOTES_SIGNALS if re.search(p, sample, re.I))

    scores = {"question_paper": qp_hits, "lab_manual": lab_hits, "lecture_notes": lec_hits}
    best = max(scores, key=scores.get)
    return best if scores[best] > 0 else "textbook"


# ─────────────────────────────────────────────────────────────────
#  Individual Processors
# ─────────────────────────────────────────────────────────────────

def process_pdf(file_path: str, subject: str) -> List[Document]:
    """Extract text from a PDF with per-page metadata."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError("PyMuPDF not installed. Run: pip install PyMuPDF")

    filename = Path(file_path).name
    doc = fitz.open(file_path)
    raw_pages: List[tuple[int, str]] = []

    for i, page in enumerate(doc):
        text = page.get_text("text").strip()
        if text:
            raw_pages.append((i + 1, text))

    doc.close()

    full_text = "\n".join(t for _, t in raw_pages)
    content_type = classify_content_type(full_text, filename)

    docs: List[Document] = []
    for page_num, text in raw_pages:
        docs.append(Document(
            page_content=text,
            metadata={
                "source": filename,
                "subject": subject,
                "content_type": content_type,
                "page": page_num,
                "file_type": "pdf",
            }
        ))
    return docs


def process_docx(file_path: str, subject: str) -> List[Document]:
    """Extract text from a DOCX file."""
    try:
        from docx import Document as DocxDocument
    except ImportError:
        raise ImportError("python-docx not installed. Run: pip install python-docx")

    filename = Path(file_path).name
    doc = DocxDocument(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    full_text = "\n".join(paragraphs)
    content_type = classify_content_type(full_text, filename)

    return [Document(
        page_content=full_text,
        metadata={
            "source": filename,
            "subject": subject,
            "content_type": content_type,
            "page": 1,
            "file_type": "docx",
        }
    )]


def process_pptx(file_path: str, subject: str) -> List[Document]:
    """Extract text from a PPTX file, one Document per slide."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")

    filename = Path(file_path).name
    prs = Presentation(file_path)
    docs: List[Document] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slide_text = "\n".join(texts)
            docs.append(Document(
                page_content=slide_text,
                metadata={
                    "source": filename,
                    "subject": subject,
                    "content_type": "lecture_notes",
                    "page": slide_num,
                    "file_type": "pptx",
                }
            ))
    return docs


# ─────────────────────────────────────────────────────────────────
#  Main Entry Point
# ─────────────────────────────────────────────────────────────────

def process_uploaded_file(file_path: str, subject: str) -> List[Document]:
    """Route to the correct processor based on file extension."""
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        raw_docs = process_pdf(file_path, subject)
    elif ext in (".docx", ".doc"):
        raw_docs = process_docx(file_path, subject)
    elif ext in (".pptx", ".ppt"):
        raw_docs = process_pptx(file_path, subject)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    # Chunk the documents
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        separators=["\n\n", "\n", ".", " "],
    )
    chunks = splitter.split_documents(raw_docs)

    # Preserve metadata after splitting
    for chunk in chunks:
        if "source" not in chunk.metadata:
            chunk.metadata["source"] = Path(file_path).name

    return chunks


def get_content_type_emoji(content_type: str) -> str:
    return {
        "question_paper": "📝",
        "lab_manual": "🔬",
        "lecture_notes": "📋",
        "textbook": "📚",
    }.get(content_type, "📄")
